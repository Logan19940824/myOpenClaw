#!/usr/bin/env python3
"""
Research & PPT Report MCP Skill (ç¾åŒ–ç‰ˆ)

æ ¹æ®ç ”ç©¶ä¸»é¢˜ï¼Œè”ç½‘æœç´¢ã€åˆ†ææ€»ç»“ï¼Œè‡ªåŠ¨ç”Ÿæˆä¸“ä¸šç¾è§‚çš„ PPT æŠ¥å‘Šã€‚

è¾“å…¥ (stdin):
{
  "topic": "<ç ”ç©¶ä¸»é¢˜>",
  "output": "<è¾“å‡ºæ–‡ä»¶è·¯å¾„>",
  "slides": "<å¹»ç¯ç‰‡æ•°é‡>",  // å¯é€‰ï¼Œé»˜è®¤ 10
  "language": "<è¯­è¨€>",       // å¯é€‰ï¼Œé»˜è®¤ "zh"
  "theme": "<ä¸»é¢˜>"          // å¯é€‰ï¼Œé»˜è®¤ "modern-blue"
}

ä¸»é¢˜é€‰é¡¹: modern-blue, business, tech, nature, gradient-purple

è¾“å‡º (stdout):
{
  "success": true,
  "output_path": "<è¾“å‡ºè·¯å¾„>",
  "slides_created": <æ•°é‡>,
  "sources": [<ä¿¡æ¯æ¥æº>],
  "error": null
}
"""

import json
import sys
import asyncio
from pathlib import Path
from datetime import datetime
from typing import List, Dict, Optional

import requests
from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE


# ä¸»é¢˜é…è‰²æ–¹æ¡ˆ
THEMES = {
    # é»˜è®¤ä¸»é¢˜
    "modern-blue": {
        "primary": RGBColor(0, 102, 204),      # è“è‰²
        "secondary": RGBColor(51, 153, 255),   # æµ…è“
        "accent": RGBColor(255, 153, 0),       # æ©™è‰²å¼ºè°ƒ
        "bg_gradient_start": RGBColor(240, 248, 255),
        "bg_gradient_end": RGBColor(220, 240, 255),
        "text_primary": RGBColor(0, 51, 102),
        "text_secondary": RGBColor(102, 102, 102),
        "font_title": "Microsoft YaHei",
        "font_body": "Microsoft YaHei",
        "bg_style": "gradient"
    },
    
    # å•†åŠ¡ä¸»é¢˜
    "business": {
        "primary": RGBColor(0, 51, 102),       # æ·±è“
        "secondary": RGBColor(51, 51, 51),     # ç°è‰²
        "accent": RGBColor(204, 153, 0),       # é‡‘è‰²
        "bg_gradient_start": RGBColor(255, 255, 255),
        "bg_gradient_end": RGBColor(245, 245, 245),
        "text_primary": RGBColor(0, 0, 0),
        "text_secondary": RGBColor(80, 80, 80),
        "font_title": "Arial",
        "font_body": "Arial",
        "bg_style": "solid"
    },
    
    # ç§‘æŠ€ä¸»é¢˜
    "tech": {
        "primary": RGBColor(16, 185, 129),     # ç»¿è‰²
        "secondary": RGBColor(99, 102, 241),   # ç´«è‰²
        "accent": RGBColor(245, 158, 11),      # æ©™è‰²
        "bg_gradient_start": RGBColor(17, 24, 39),
        "bg_gradient_end": RGBColor(31, 41, 55),
        "text_primary": RGBColor(255, 255, 255),
        "text_secondary": RGBColor(200, 200, 200),
        "font_title": "Segoe UI",
        "font_body": "Segoe UI",
        "bg_style": "dark"
    },
    
    # è‡ªç„¶ä¸»é¢˜
    "nature": {
        "primary": RGBColor(34, 139, 34),      # æ£®æ—ç»¿
        "secondary": RGBColor(85, 107, 47),    # æ©„æ¦„è‰²
        "accent": RGBColor(255, 165, 0),       # æ©™è‰²
        "bg_gradient_start": RGBColor(240, 255, 240),
        "bg_gradient_end": RGBColor(220, 255, 220),
        "text_primary": RGBColor(0, 100, 0),
        "text_secondary": RGBColor(60, 80, 60),
        "font_title": "Georgia",
        "font_body": "Calibri",
        "bg_style": "gradient"
    },
    
    # æ¸å˜ç´«ä¸»é¢˜
    "gradient-purple": {
        "primary": RGBColor(128, 0, 128),      # ç´«è‰²
        "secondary": RGBColor(255, 0, 255),    # ç²‰ç´«
        "accent": RGBColor(0, 255, 255),       # é’è‰²
        "bg_gradient_start": RGBColor(240, 230, 255),
        "bg_gradient_end": RGBColor(255, 240, 255),
        "text_primary": RGBColor(64, 0, 64),
        "text_secondary": RGBColor(100, 100, 100),
        "font_title": "Verdana",
        "font_body": "Verdana",
        "bg_style": "gradient"
    },
    
    # ğŸ¨ ç²¾å“æ¨¡æ¿1: æ¸å˜æ©™è‰² (æ´»åŠ›é£æ ¼)
    "gradient-orange": {
        "name": "æ´»åŠ›æ©™",
        "primary": RGBColor(255, 140, 0),      # æ©™è‰²
        "secondary": RGBColor(255, 69, 0),     # çº¢æ©™
        "accent": RGBColor(255, 215, 0),       # é‡‘è‰²
        "bg_gradient_start": RGBColor(255, 248, 240),
        "bg_gradient_end": RGBColor(255, 224, 178),
        "text_primary": RGBColor(139, 69, 19),
        "text_secondary": RGBColor(160, 82, 45),
        "font_title": "Microsoft YaHei",
        "font_body": "Microsoft YaHei",
        "bg_style": "gradient",
        "has_watermark": True
    },
    
    # ğŸ¨ ç²¾å“æ¨¡æ¿2: é«˜çº§é»‘é‡‘ (å¥¢åé£æ ¼)
    "premium-black-gold": {
        "name": "é«˜çº§é»‘é‡‘",
        "primary": RGBColor(218, 165, 32),     # é‡‘è‰²
        "secondary": RGBColor(139, 69, 19),    # æ£•è‰²
        "accent": RGBColor(255, 215, 0),       # é‡‘è‰²
        "bg_gradient_start": RGBColor(30, 30, 30),
        "bg_gradient_end": RGBColor(50, 50, 50),
        "text_primary": RGBColor(255, 215, 0),
        "text_secondary": RGBColor(200, 200, 200),
        "font_title": "Arial Black",
        "font_body": "Georgia",
        "bg_style": "dark",
        "has_watermark": True,
        "decorative_lines": True
    },
    
    # ğŸ¨ ç²¾å“æ¨¡æ¿3: æç®€ç™½ (å•†åŠ¡æç®€)
    "minimal-white": {
        "name": "æç®€ç™½",
        "primary": RGBColor(0, 0, 0),          # é»‘è‰²
        "secondary": RGBColor(128, 128, 128),  # ç°è‰²
        "accent": RGBColor(0, 0, 0),           # é»‘è‰²
        "bg_gradient_start": RGBColor(255, 255, 255),
        "bg_gradient_end": RGBColor(255, 255, 255),
        "text_primary": RGBColor(0, 0, 0),
        "text_secondary": RGBColor(80, 80, 80),
        "font_title": "Helvetica",
        "font_body": "Helvetica",
        "bg_style": "solid",
        "has_watermark": False,
        "decorative_lines": False
    },
    
    # ğŸ¨ ç²¾å“æ¨¡æ¿4: æ¸å˜é’è“ (ç§‘æŠ€æœªæ¥)
    "tech-future": {
        "name": "ç§‘æŠ€æœªæ¥",
        "primary": RGBColor(0, 206, 209),      # æ·±é’è‰²
        "secondary": RGBColor(30, 144, 255),   # é“å¥‡è“
        "accent": RGBColor(0, 255, 127),       # æ˜¥ç»¿è‰²
        "bg_gradient_start": RGBColor(0, 30, 60),
        "bg_gradient_end": RGBColor(0, 60, 100),
        "text_primary": RGBColor(255, 255, 255),
        "text_secondary": RGBColor(180, 220, 255),
        "font_title": "Segoe UI",
        "font_body": "Segoe UI",
        "bg_style": "dark",
        "has_watermark": True,
        "decorative_circuits": True
    },
    
    # ğŸ¨ ç²¾å“æ¨¡æ¿5: çº¢è‰²ä¸­å›½é£ (å–œåº†é£æ ¼)
    "chinese-red": {
        "name": "ä¸­å›½çº¢",
        "primary": RGBColor(178, 34, 34),      # æ·±çº¢
        "secondary": RGBColor(220, 20, 60),    # çŒ©çº¢
        "accent": RGBColor(255, 215, 0),       # é‡‘è‰²
        "bg_gradient_start": RGBColor(255, 240, 240),
        "bg_gradient_end": RGBColor(255, 200, 200),
        "text_primary": RGBColor(139, 0, 0),
        "text_secondary": RGBColor(178, 34, 34),
        "font_title": "Microsoft YaHei",
        "font_body": "Microsoft YaHei",
        "bg_style": "gradient",
        "has_watermark": False,
        "decorative_pattern": "cloud"
    }
}

# é»˜è®¤é…ç½®
DEFAULT_SLIDES = 10
DEFAULT_THEME = "modern-blue"

LANGUAGE_MAP = {
    "zh": {"title": "ç ”ç©¶æŠ¥å‘Š", "toc": "ç›®å½•", "summary": "æ€»ç»“", "sources": "å‚è€ƒèµ„æ–™"},
    "en": {"title": "Research Report", "toc": "Contents", "summary": "Summary", "sources": "References"}
}


class BeautifulPPTSkill:
    """ç¾è§‚ PPT ç”Ÿæˆå™¨"""
    
    def __init__(self):
        self.session = requests.Session()
        self.session.headers.update({
            "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36"
        })
    
    def search_web(self, query: str, num_results: int = 5) -> List[Dict]:
        """è”ç½‘æœç´¢"""
        try:
            url = f"https://duckduckgo.com/html/?q={query}&kl=us-en&ia=web"
            response = self.session.get(url, timeout=10)
            response.raise_for_status()
            
            soup = BeautifulSoup(response.text, 'html.parser')
            results = []
            
            for result in soup.find_all('div', class_='result')[:num_results]:
                try:
                    title_elem = result.find('a', class_='result__a')
                    snippet_elem = result.find('a', class_='result__snippet')
                    
                    if title_elem:
                        results.append({
                            "title": title_elem.get_text(strip=True),
                            "url": title_elem.get('href', ''),
                            "snippet": snippet_elem.get_text(strip=True) if snippet_elem else ""
                        })
                except Exception:
                    continue
            
            return results[:num_results]
            
        except Exception as e:
            print(f"æœç´¢å‡ºé”™: {e}", file=sys.stderr)
            return []
    
    def analyze_info(self, topic: str, sources: List[Dict], language: str = "zh") -> Dict:
        """åˆ†æç”Ÿæˆ PPT å†…å®¹å¤§çº²"""
        lang = LANGUAGE_MAP.get(language, LANGUAGE_MAP["zh"])
        
        content = {
            "title": topic,
            "subtitle": f"{lang['title']} | {datetime.now().strftime('%Yå¹´%mæœˆ')}",
            "slides": []
        }
        
        # ä¸“ä¸šçš„ PPT ç»“æ„
        slide_templates = [
            ("ç›®å½•", ["ç ”ç©¶èƒŒæ™¯", "å¸‚åœºåˆ†æ", "å‘å±•è¶‹åŠ¿", "æŒ‘æˆ˜ä¸æœºé‡", "æœªæ¥å±•æœ›"]),
            ("ç ”ç©¶èƒŒæ™¯", [f"ç ”ç©¶ä¸»é¢˜ï¼š{topic}", "ç ”ç©¶ç›®çš„ä¸æ„ä¹‰", "ç ”ç©¶æ–¹æ³•è¯´æ˜", "æ•°æ®æ¥æºä»‹ç»"]),
            ("å¸‚åœºåˆ†æ", ["å…¨çƒå¸‚åœºè§„æ¨¡", "ä¸­å›½å¸‚åœºç°çŠ¶", "ä¸»è¦å‚ä¸è€…", "å¢é•¿é©±åŠ¨å› ç´ "]),
            ("å‘å±•è¶‹åŠ¿", ["æŠ€æœ¯åˆ›æ–°æ–¹å‘", "å•†ä¸šæ¨¡å¼æ¼”è¿›", "æ”¿ç­–ç¯å¢ƒå½±å“", "æ¶ˆè´¹è€…éœ€æ±‚å˜åŒ–"]),
            ("ç«äº‰æ ¼å±€", ["ä¸»è¦ç«äº‰å¯¹æ‰‹", "å¸‚åœºä»½é¢åˆ†å¸ƒ", "ç«äº‰ä¼˜åŠ¿å¯¹æ¯”", "å¸‚åœºè¿›å…¥å£å’"]),
            ("æŒ‘æˆ˜ä¸æœºé‡", ["è¡Œä¸šå‘å±•ç—›ç‚¹", "æ½œåœ¨å¢é•¿ç©ºé—´", "é£é™©å› ç´ è¯†åˆ«", "åº”å¯¹ç­–ç•¥å»ºè®®"]),
            ("æœªæ¥å±•æœ›", ["çŸ­æœŸé¢„æµ‹ (1-2å¹´)", "ä¸­æœŸé¢„æµ‹ (3-5å¹´)", "é•¿æœŸè¶‹åŠ¿åˆ¤æ–­", "æˆ˜ç•¥å»ºè®®"]),
            (lang['summary'], ["æ ¸å¿ƒå‘ç°æ€»ç»“", "å…³é”®æ•°æ®æ”¯æ’‘", "å†³ç­–å»ºè®®", "åç»­ç ”ç©¶æ–¹å‘"]),
            (lang['sources'], ["ä¿¡æ¯æ¥æºè¯´æ˜", "ç ”ç©¶å±€é™æ€§", "å…è´£å£°æ˜"])
        ]
        
        for slide_title, bullets in slide_templates:
            content["slides"].append({
                "title": slide_title,
                "bullets": bullets
            })
        
        content["sources"] = [s["url"] for s in sources[:5]] if sources else ["å…¬å¼€ä¿¡æ¯æ•´ç†"]
        
        return content
    
    def add_background(self, slide, theme: Dict, is_cover: bool = False):
        """æ·»åŠ èƒŒæ™¯"""
        if is_cover:
            # å°é¢èƒŒæ™¯
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, 0, 0, 
                Inches(13.333), Inches(7.5)
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = theme["bg_gradient_start"]
            shape.line.fill.background()
            
            # å°é¢åº•éƒ¨è£…é¥°
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, 0, Inches(5), 
                Inches(13.333), Inches(2.5)
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = theme["primary"]
            shape.fill.transparency = 0.1
            shape.line.fill.background()
        else:
            # å†…å®¹é¡µé¡¶éƒ¨è£…é¥°æ¡
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, 0, 0, 
                Inches(13.333), Inches(0.15)
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = theme["primary"]
            shape.line.fill.background()
            
            # åº•éƒ¨è£…é¥°æ¡
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, 0, Inches(7.35), 
                Inches(13.333), Inches(0.15)
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = theme["secondary"]
            shape.line.fill.background()
    
    def add_page_number(self, slide, page_num: int, theme: Dict, total: int):
        """æ·»åŠ é¡µç """
        textbox = slide.shapes.add_textbox(
            Inches(12), Inches(7), Inches(1), Inches(0.3)
        )
        tf = textbox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"{page_num}/{total}"
        p.font.size = Pt(12)
        p.font.color.rgb = theme["text_secondary"]
        p.alignment = PP_ALIGN.RIGHT
    
    def add_decorative_elements(self, slide, theme: Dict, position: str = "corner"):
        """æ·»åŠ è£…é¥°å…ƒç´ """
        # å³ä¸‹è§’è£…é¥°åœ†åœˆ
        shape = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, 
            Inches(11.5), Inches(5.5), 
            Inches(1.5), Inches(1.5)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = theme["secondary"]
        shape.fill.transparency = 0.7
        shape.line.fill.background()
    
    def create_ppt(self, content: Dict, output_path: str, theme_name: str = DEFAULT_THEME) -> Dict:
        """ç”Ÿæˆç¾è§‚ PPT"""
        try:
            theme = THEMES.get(theme_name, THEMES[DEFAULT_THEME])
            prs = Presentation()
            
            # è®¾ç½®é¡µé¢å¤§å° (16:9)
            prs.slide_width = Inches(13.333)
            prs.slide_height = Inches(7.5)
            
            total_slides = len(content["slides"])
            
            for i, slide_data in enumerate(content["slides"]):
                if i == 0:
                    # å°é¢é¡µ
                    slide = prs.slides.add_slide(prs.slide_layouts[6])
                    self.add_background(slide, theme, is_cover=True)
                    
                    # è£…é¥°åœ†åœˆ
                    shape = slide.shapes.add_shape(
                        MSO_SHAPE.OVAL, Inches(9.5), Inches(1), 
                        Inches(3.5), Inches(3.5)
                    )
                    shape.fill.solid()
                    shape.fill.fore_color.rgb = theme["secondary"]
                    shape.fill.transparency = 0.2
                    shape.line.fill.background()
                    
                    # ä¸»æ ‡é¢˜
                    title_box = slide.shapes.add_textbox(
                        Inches(0.8), Inches(2), Inches(11.733), Inches(1.5)
                    )
                    tf = title_box.text_frame
                    tf.word_wrap = True
                    p = tf.paragraphs[0]
                    p.text = slide_data["title"]
                    p.font.size = Pt(44)
                    p.font.bold = True
                    p.font.color.rgb = theme["text_primary"]
                    p.alignment = PP_ALIGN.LEFT
                    
                    # å‰¯æ ‡é¢˜
                    subtitle_box = slide.shapes.add_textbox(
                        Inches(0.8), Inches(3.8), Inches(11.733), Inches(0.8)
                    )
                    tf = subtitle_box.text_frame
                    p = tf.paragraphs[0]
                    p.text = content["subtitle"]
                    p.font.size = Pt(20)
                    p.font.color.rgb = theme["secondary"]
                    p.alignment = PP_ALIGN.LEFT
                    
                    # åº•éƒ¨ä¿¡æ¯
                    footer_box = slide.shapes.add_textbox(
                        Inches(0.8), Inches(6.2), Inches(11.733), Inches(0.5)
                    )
                    tf = footer_box.text_frame
                    p = tf.paragraphs[0]
                    p.text = "æŒ‰ Enter é”®ç»§ç»­ | Press Enter to continue"
                    p.font.size = Pt(12)
                    p.font.color.rgb = theme["text_secondary"]
                    p.alignment = PP_ALIGN.LEFT
                    
                else:
                    # å†…å®¹é¡µ
                    slide = prs.slides.add_slide(prs.slide_layouts[6])
                    self.add_background(slide, theme)
                    
                    # å·¦ä¾§è£…é¥°æ¡
                    shape = slide.shapes.add_shape(
                        MSO_SHAPE.RECTANGLE, 
                        Inches(0.3), Inches(0.8), 
                        Inches(0.08), Inches(5.5)
                    )
                    shape.fill.solid()
                    shape.fill.fore_color.rgb = theme["primary"]
                    shape.line.fill.background()
                    
                    # æ ‡é¢˜åŒºåŸŸ
                    title_box = slide.shapes.add_textbox(
                        Inches(0.6), Inches(0.4), 
                        Inches(12), Inches(0.7)
                    )
                    tf = title_box.text_frame
                    p = tf.paragraphs[0]
                    p.text = f"0{i}. {slide_data['title']}"
                    p.font.size = Pt(28)
                    p.font.bold = True
                    p.font.color.rgb = theme["primary"]
                    
                    # å†…å®¹åŒºåŸŸ
                    content_box = slide.shapes.add_textbox(
                        Inches(0.8), Inches(1.5), 
                        Inches(12), Inches(5.3)
                    )
                    tf = content_box.text_frame
                    tf.word_wrap = True
                    
                    for j, bullet in enumerate(slide_data["bullets"]):
                        if j == 0:
                            p = tf.paragraphs[0]
                        else:
                            p = tf.add_paragraph()
                        p.text = f"âœ“ {bullet}"
                        p.font.size = Pt(18)
                        p.font.color.rgb = theme["text_primary"]
                        p.space_before = Pt(14)
                        
                        # ç¬¬ä¸€ä¸ªè¦ç‚¹åŠ å¤§åŠ ç²—
                        if j == 0:
                            p.font.size = Pt(20)
                            p.font.bold = True
                    
                    # è£…é¥°å…ƒç´ 
                    self.add_decorative_elements(slide, theme)
                    
                    # é¡µç 
                    self.add_page_number(slide, i, theme, total_slides - 1)
            
            # ä¿å­˜æ–‡ä»¶
            output_file = Path(output_path)
            output_file.parent.mkdir(parents=True, exist_ok=True)
            prs.save(str(output_file))
            
            return {
                "success": True,
                "output_path": str(output_file),
                "slides_created": len(content["slides"]),
                "theme": theme_name,
                "sources": content.get("sources", []),
                "error": None
            }
            
        except Exception as e:
            print(f"ç”Ÿæˆ PPT é”™è¯¯: {e}", file=sys.stderr)
            return {
                "success": False,
                "output_path": None,
                "slides_created": 0,
                "sources": [],
                "error": str(e)
            }
    
    async def execute(self, topic: str, output: str, slides: int = DEFAULT_SLIDES, 
                     language: str = "zh", theme: str = DEFAULT_THEME) -> Dict:
        """æ‰§è¡Œå®Œæ•´çš„ PPT ç”Ÿæˆæµç¨‹"""
        print(f"æ­£åœ¨ç ”ç©¶ä¸»é¢˜: {topic}", file=sys.stderr)
        
        # 1. è”ç½‘æœç´¢
        print("æ­¥éª¤ 1/3: è”ç½‘æœç´¢...", file=sys.stderr)
        sources = self.search_web(topic, num_results=5)
        
        if not sources:
            sources = [{"title": f"å…³äº {topic} çš„ç ”ç©¶æŠ¥å‘Š", "url": "ç½‘ç»œèµ„æº"}]
        
        # 2. åˆ†ææ€»ç»“
        print("æ­¥éª¤ 2/3: åˆ†ææ€»ç»“...", file=sys.stderr)
        content = self.analyze_info(topic, sources, language)
        
        # 3. ç”Ÿæˆ PPT
        print(f"æ­¥éª¤ 3/3: ç”Ÿæˆ PPT (ä¸»é¢˜: {theme})...", file=sys.stderr)
        result = self.create_ppt(content, output, theme)
        
        return result


def main():
    """ä¸»å‡½æ•°"""
    input_data = sys.stdin.read().strip()
    
    if not input_data:
        result = {"success": False, "output_path": None, "slides_created": 0, "error": "æœªæ”¶åˆ°è¾“å…¥æ•°æ®"}
        print(json.dumps(result))
        return
    
    try:
        data = json.loads(input_data)
        topic = data.get("topic")
        output = data.get("output")
        slides = data.get("slides", DEFAULT_SLIDES)
        language = data.get("language", "zh")
        theme = data.get("theme", DEFAULT_THEME)
        
        if not topic or not output:
            result = {"success": False, "output_path": None, "slides_created": 0, "error": "ç¼ºå°‘å¿…è¦å‚æ•°"}
            print(json.dumps(result))
            return
        
        skill = BeautifulPPTSkill()
        result = asyncio.run(skill.execute(topic, output, slides, language, theme))
        print(json.dumps(result))
        
    except Exception as e:
        result = {"success": False, "output_path": None, "slides_created": 0, "error": str(e)}
        print(json.dumps(result))


if __name__ == "__main__":
    main()
